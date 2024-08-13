from __future__ import annotations
import os
import json
from typing import List, Optional
import pptx
from pptx import Presentation
import pptx.presentation
from pptx.util import Inches, Pt
from pptx.slide import Slide as PptxSlide
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.text.text import TextFrame

class Slide:
    def __init__(self, title: str, content: str, image_path: Optional[str] = None):
        self.title = title
        self.content = content
        self.image_path = image_path

    def to_json(self) -> str:
        return json.dumps({
            "title": self.title,
            "content": self.content,
            "image_path": self.image_path
        }).replace("\\n", "\n")

    @classmethod
    def from_dict(cls, data: dict[str,str]) -> Slide:
        return cls(
            title=data['title'],
            content=data['content'],
            image_path=data.get('image_path')
        )

class PptxPresentation:
    def __init__(self, title: str, subtitle: str, slides: List[Slide]):
        self.prs: pptx.presentation.Presentation = Presentation()  # Using Any to avoid mypy error
        self.title: str = title
        self.subtitle: str = subtitle
        self.slides: List[Slide] = slides
        self._create_title_slide()
        self._create_content_slides()

    def _create_title_slide(self) -> None:
        layout = self.prs.slide_layouts[0]
        slide: PptxSlide = self.prs.slides.add_slide(layout)
        title_shape: Shape = slide.shapes.title
        subtitle_shape: Shape = slide.placeholders[1]
        title_shape.text = self.title
        subtitle_shape.text = self.subtitle

    def _create_content_slides(self) -> None:
        for slide in self.slides:
            self._add_content_slide(slide)

    def _add_content_slide(self, slide: Slide) -> PptxSlide:
        layout = self.prs.slide_layouts[1]
        pptx_slide: PptxSlide = self.prs.slides.add_slide(layout)
        title_shape: Shape = pptx_slide.shapes.title
        content_shape: Shape = pptx_slide.placeholders[1]
        title_shape.text = slide.title
        tf: TextFrame = content_shape.text_frame
        tf.clear()  # Clear existing text
        slide_contents = slide.content.split('\n')
        for slide_content in slide_contents:
            slide_content = slide_content.strip("• ").strip()
            p = tf.add_paragraph()
            p.text = slide_content
            p.font.size = Pt(18)
            if slide.image_path:
                self._add_image_to_slide(pptx_slide, slide.image_path)
        return pptx_slide

    def _add_image_to_slide(self, pptx_slide: PptxSlide, image_path: str) -> None:
        content_shape: Shape = pptx_slide.placeholders[1]
        content_shape.width = Inches(5)
        img_left = Inches(7)
        img_top = Inches(2)
        img_width = Inches(3)
        pptx_slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)

    def save(self) -> None:
        self.prs.save(f'{self.title.replace(" ", "_")}.pptx')

    def to_json(self) -> str:
        return json.dumps({
            "title": self.title,
            "subtitle": self.subtitle,
            "slides": [slide.to_json() for slide in self.slides]
        })

    @classmethod
    def from_json(cls, json_data: str) -> PptxPresentation:
        data = json.loads(json_data)
        slides = [Slide.from_dict(slide_data) for slide_data in data['slides']]
        return cls(
            title=data['title'],
            subtitle=data['subtitle'],
            slides=slides
        )
        
if __name__ == "__main__":
    slides_1 = [
        Slide("Hybrid Approach", "• Combination of rule-based logic from 'marvin' with ML models\n• Enables gradual integration and maintains a fallback system"),
        Slide("Modular ML Integration", "• Development of modular ML components for specific strategy areas\n• Facilitates testing, comparisons, and incremental improvements"),
        Slide("Real-time Adaptation", "• Implementation of online learning for fine-tuning during games\n• Helps adapt to opponent strategies in real-time"),
        Slide("Transfer Learning", "• Utilization of pre-trained models from similar domains\n• Reduces training time and improves initial performance"),
        Slide("Opponent Modeling", "• Implementation of a system for predicting and countering various strategies\n• Particularly useful in tournament environments"),
        Slide("Hierarchical Learning", "• Separation of high-level strategy decisions and low-level execution")
    ]
    presentation_1 = PptxPresentation("ER-Force Strategy Optimization", "Strategy meeting 2024", slides_1)
    presentation_1.save()