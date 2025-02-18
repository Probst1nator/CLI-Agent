import pyglet
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

@dataclass
class SlideContent:
    """Represents the content of a single slide."""
    title: str
    bullets: List[str]
    background_color: Tuple[int, int, int] = (255, 255, 255)

class Slide:
    """Manages individual slide creation and rendering."""
    
    def __init__(self, content: SlideContent, window_size: Tuple[int, int]):
        """Initialize a slide with content and window dimensions.
        
        Args:
            content: SlideContent object containing slide information
            window_size: Tuple of window width and height
        """
        self.content = content
        self.window_size = window_size
        self._setup_text()

    def _setup_text(self) -> None:
        """Setup text elements for the slide."""
        try:
            self.title_label = pyglet.text.Label(
                self.content.title,
                font_name='Arial',
                font_size=36,
                x=self.window_size[0]//2,
                y=self.window_size[1]-100,
                anchor_x='center',
                anchor_y='center'
            )
            
            self.bullet_labels = []
            for i, bullet in enumerate(self.content.bullets):
                y_pos = self.window_size[1] - 200 - (i * 50)
                label = pyglet.text.Label(
                    f"• {bullet}",
                    font_name='Arial',
                    font_size=24,
                    x=100,
                    y=y_pos,
                    anchor_x='left',
                    anchor_y='center'
                )
                self.bullet_labels.append(label)
        except Exception as e:
            raise RuntimeError(f"Failed to setup slide text: {str(e)}")

    def draw(self) -> None:
        """Draw the slide content."""
        try:
            # Draw background
            pyglet.graphics.draw(
                4, pyglet.gl.GL_QUADS,
                ('v2f', (0, 0, self.window_size[0], 0,
                        self.window_size[0], self.window_size[1],
                        0, self.window_size[1])),
                ('c3B', self.content.background_color * 4)
            )
            
            # Draw text elements
            self.title_label.draw()
            for label in self.bullet_labels:
                label.draw()
        except Exception as e:
            raise RuntimeError(f"Failed to draw slide: {str(e)}")

class Presentation:
    """Manages the presentation window and slides."""
    
    def __init__(self, title: str, slides_content: List[SlideContent]):
        """Initialize the presentation.
        
        Args:
            title: Presentation title
            slides_content: List of SlideContent objects
        """
        self.window_size = (1024, 768)
        self.current_slide = 0
        self.window = pyglet.window.Window(*self.window_size, caption=title)
        self.slides = self._create_slides(slides_content)
        self._setup_event_handlers()

    def _create_slides(self, slides_content: List[SlideContent]) -> List[Slide]:
        """Create slide objects from content."""
        return [Slide(content, self.window_size) for content in slides_content]
    
    def _setup_event_handlers(self):
        @self.window.event
        def on_draw():
            self.window.clear()
            if len(self.slides) > self.current_slide:
                self.slides[self.current_slide].draw()

        @self.window.event
        def on_key_press(symbol, modifiers):
            if symbol == pyglet.window.key.RIGHT:
                self.current_slide = (self.current_slide + 1) % len(self.slides)
            elif symbol == pyglet.window.key.LEFT:
                self.current_slide = (self.current_slide - 1) % len(self.slides)

    def run(self):
        pyglet.app.run()

def create_ai_presentation() -> None:
    """Create and run the AI advancements presentation."""
    try:
        slides_content = [
            SlideContent(
                "AI Advancements 2023",
                ["Latest developments in artificial intelligence",
                 "Impact on various industries",
                 "Future predictions"]
            ),
            SlideContent(
                "Key Breakthroughs",
                ["Large Language Models",
                 "Computer Vision Advances",
                 "Robotics Integration"],
                (240, 240, 255)
            ),
            SlideContent(
                "Future Outlook",
                ["Increased adoption",
                 "Ethical considerations",
                 "New applications"],
                (255, 240, 240)
            )
        ]
        
        presentation = Presentation()
        for content in slides_content:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = content.title
            bullet_points = slide.placeholders[1]
            for i, bullet in enumerate(content.bullets):
                p = bullet_points.add_paragraph()
                p.level = 0
                p.text = f"{i+1}. {bullet}"
        
        presentation.save("AI_Advancements.pptx")
        
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    create_ai_presentation()