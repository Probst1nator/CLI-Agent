from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import sys

class JNJCorruptionPresentation:
    """Class to create and manage a PowerPoint presentation about J&J corruption cases."""

    def __init__(self) -> None:
        """Initialize presentation with basic settings."""
        self.prs = Presentation()
        self.findings: List[Dict[str, str]] = [
            {
                "title": "Foreign Corrupt Practices Act Settlement",
                "content": "J&J agreed to pay $70 million to resolve FCPA violations",
                "date": "2011",
                "source": "U.S. Department of Justice"
            },
            {
                "title": "Healthcare Fraud Settlement",
                "content": "Settlement of $2.2 billion for promoting drugs for unapproved uses",
                "date": "2013",
                "source": "DOJ Press Release"
            }
        ]

    def create_title_slide(self) -> None:
        """Create the presentation title slide."""
        try:
            layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Johnson & Johnson: Corruption Cases Analysis"
            subtitle.text = f"Period: 2011-2023\nCompiled: {datetime.now().strftime('%B %d, %Y')}"

            # Format title
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        except Exception as e:
            print(f"Error creating title slide: {str(e)}")
            raise

    def create_content_slide(self, finding: Dict[str, str]) -> None:
        """Create a content slide for each finding.

        Args:
            finding: Dictionary containing finding details
        """
        try:
            layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(layout)
            
            # Add and format title
            title = slide.shapes.title
            title.text = finding["title"]
            title.text_frame.paragraphs[0].font.size = Pt(36)
            
            # Add content
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = finding["content"]
            
            # Add source and date
            p = tf.add_paragraph()
            p.text = f"\nDate: {finding['date']}"
            p = tf.add_paragraph()
            p.text = f"Source: {finding['source']}"
            
            # Format content
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.alignment = PP_ALIGN.LEFT
        except Exception as e:
            print(f"Error creating content slide: {str(e)}")
            raise

    def create_presentation(self) -> None:
        """Create the complete presentation with all slides."""
        try:
            # Create title slide
            self.create_title_slide()
            
            # Create content slides
            for finding in self.findings:
                self.create_content_slide(finding)
                
        except Exception as e:
            print(f"Error creating presentation: {str(e)}")
            raise

    def save_presentation(self, filename: str = "corruption_in_jnj.pptx") -> None:
        """Save the presentation to a file.

        Args:
            filename: Name of the file to save the presentation
        """
        try:
            self.prs.save(filename)
            print(f"Presentation saved successfully as {filename}")
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")
            raise

def main() -> None:
    """Main function to create and save the presentation."""
    try:
        # Create presentation instance
        presentation = JNJCorruptionPresentation()
        
        # Create and save presentation
        presentation.create_presentation()
        presentation.save_presentation()
        
    except Exception as e:
        print(f"An error occurred: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()